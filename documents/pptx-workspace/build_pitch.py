#!/usr/bin/env python3
"""Build the Secure Insights AI pitch deck as a real .pptx (16:9, dark theme).

    python3 build_pitch.py            # -> ../secure-insights-pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- palette ----
BG     = RGBColor(0x0A, 0x0B, 0x10)
PANEL  = RGBColor(0x16, 0x19, 0x22)
LINE   = RGBColor(0x2A, 0x2F, 0x3A)
INK    = RGBColor(0xEE, 0xF1, 0xF7)
MUTED  = RGBColor(0x8B, 0x93, 0xA5)
RED    = RGBColor(0xE4, 0x00, 0x2B)
CYAN   = RGBColor(0x2E, 0xE6, 0xD6)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
AMBER  = RGBColor(0xF7, 0xB3, 0x2B)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)

DISP = "Segoe UI"      # display / body
MONO = "Consolas"      # eyebrows / labels / data

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

MX = 0.9               # left margin
CW = 13.333 - 2 * MX   # content width


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _set(tf, runs, size, color, bold=False, font=DISP, align=PP_ALIGN.LEFT,
         spacing=None, line=1.08):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for txt, col, bd in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.name = font
        r.font.bold = bd; r.font.color.rgb = col
        if spacing is not None:
            _letter_spacing(r, spacing)
    return p


def _letter_spacing(run, pts):
    # crude tracking via XML spc attribute (100 = 1pt)
    run.font._rPr.set("spc", str(int(pts * 100)))


def box(s, x, y, w, h, runs, size, color, bold=False, font=DISP,
        align=PP_ALIGN.LEFT, spacing=None, anchor=MSO_ANCHOR.TOP, line=1.08):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _set(tf, runs, size, color, bold, font, align, spacing, line)
    return tb


def para(tb, runs, size, color, bold=False, font=DISP, spacing=None, line=1.35, space_before=6):
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.line_spacing = line
    p.space_before = Pt(space_before)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for txt, col, bd in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.name = font
        r.font.bold = bd; r.font.color.rgb = col
    return p


def rect(s, x, y, w, h, fill, line_col=None, line_w=0.75, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_col is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_col; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def eyebrow(s, text):
    box(s, MX, 1.15, CW, 0.35, text.upper(), 12.5, CYAN, bold=True, font=MONO, spacing=2.2)


def title(s, runs, size=34):
    box(s, MX, 1.62, CW, 1.5, runs, size, INK, bold=True, spacing=-0.3, line=1.02)


def lead(s, runs, y=3.05, w=None, size=17):
    box(s, MX, y, w or 9.6, 1.6, runs, size, RGBColor(0xC9, 0xCF, 0xDB), line=1.4)


def card(s, x, y, w, h, accent, heading, body, num=None):
    rect(s, x, y, w, h, PANEL, LINE)
    rect(s, x, y, 0.06, h, accent, rounded=False)
    ty = y + 0.2
    if num:
        box(s, x + 0.26, ty, w - 0.5, 0.28, num.upper(), 11, accent, bold=True, font=MONO, spacing=1.2)
        ty += 0.32
    tb = box(s, x + 0.26, ty, w - 0.5, h - (ty - y) - 0.15, heading, 16.5, INK, bold=True, line=1.08)
    para(tb, body, 13, MUTED, line=1.34, space_before=7)


def pill(s, x, y, text, color=MUTED, border=LINE, dot=False):
    w = 0.16 + 0.092 * len(text)
    shp = rect(s, x, y, w, 0.34, PANEL, border, rounded=True)
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    dx = x + 0.16
    if dot:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.14), Inches(y + 0.12), Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = GREEN; d.line.fill.background(); d.shadow.inherit = False
        dx = x + 0.32
    box(s, dx, y + 0.02, w, 0.3, text, 10.5, color, bold=True, font=MONO, spacing=0.6, anchor=MSO_ANCHOR.MIDDLE)
    return x + w + 0.16


def footer(s):
    box(s, MX, 7.05, 6, 0.3, "■ SECURE INSIGHTS AI", 10.5, RED, bold=True, font=MONO, spacing=1.2)


# ============================================================ SLIDES

# 1 · TITLE
s = slide()
eyebrow(s, "Agentic AI · Enterprise Operations")
box(s, MX, 1.7, 11.5, 2.4, [("See it in ", INK, True), ("minutes", CYAN, True), (".", INK, True)],
    52, INK, bold=True, spacing=-0.6, line=1.0)
box(s, MX, 2.75, 11.5, 1.4, [("Fix it ", INK, True), ("same shift", RED, True), (".", INK, True)],
    52, INK, bold=True, spacing=-0.6, line=1.0)
lead(s, "An autonomous platform that forecasts demand, detects what's breaking across 250+ "
        "restaurants, explains why, and deploys the fix — before the shift ends.", y=4.15, w=10.6, size=18)
x = pill(s, MX, 5.7, "Live · forecaster.secureinsights.ai", GREEN, RGBColor(0x1e, 0x4d, 0x3a), dot=True)
x = pill(s, x, 5.7, "Secure Insights AI")
pill(s, x, 5.7, "KFC Vietnam")

# 2 · PROBLEM
s = slide()
eyebrow(s, "The problem")
title(s, "You're always the last to know.", 40)
lead(s, [("KFC Vietnam runs ", RGBColor(0xC9,0xCF,0xDB), False), ("250+ restaurants", INK, True),
         (" and millions of daily transactions across dine-in, kiosk, delivery, and app. When something "
          "breaks — a register outage, a delivery-partner blackout, comp fraud — Ops & Finance only find "
          "out from a ", RGBColor(0xC9,0xCF,0xDB), False), ("manual report review…", INK, True)], y=2.95, w=11)
box(s, MX, 4.35, 8, 1.4, "24–48h", 78, RED, bold=True, spacing=-1.5)
box(s, MX, 5.75, 10, 0.4, "AFTER THE MONEY ALREADY WALKED OUT THE DOOR", 12, MUTED, bold=True, font=MONO, spacing=1.2)
box(s, MX, 6.3, 11, 0.6, [("And there's ", MUTED, False), ("no predictive baseline", CYAN, True),
    (" — no way to tell a real problem from an ordinary rainy Tuesday after a holiday.", MUTED, False)], 15, MUTED, line=1.4)

# 3 · WHY IT MATTERS
s = slide()
eyebrow(s, "Why it matters")
title(s, [("In operations, latency ", INK, True), ("is", RED, True), (" loss.", INK, True)], 40)
lead(s, [("The damage of an incident is the ", RGBColor(0xC9,0xCF,0xDB), False),
         ("rate of harm × the time you stay blind to it.", INK, True),
         (" A store closed since lunch, found tomorrow afternoon, is a full day of revenue you can never "
          "get back.", RGBColor(0xC9,0xCF,0xDB), False)], y=2.95, w=11)
cy, cw, ch = 4.55, (CW - 2*0.28)/3, 1.9
card(s, MX,               cy, cw, ch, RED, "It compounds hourly", "A POS outage, a fraud pattern, a stockout — each keeps bleeding every shift until someone notices.")
card(s, MX+cw+0.28,       cy, cw, ch, RED, "It doesn't scale", "More stores → more reports → longer lag. Manual review gets worse as the chain grows.")
card(s, MX+2*(cw+0.28),   cy, cw, ch, RED, "It burns your best people", "Analysts spend their days finding problems instead of fixing them.")

# 4 · SOLUTION
s = slide()
eyebrow(s, "The solution")
title(s, "Not a dashboard. An operator that acts.", 38)
lead(s, [("A ", RGBColor(0xC9,0xCF,0xDB), False), ("closed-loop agent", INK, True),
         (" that watches every store, reasons about what's wrong, and ", RGBColor(0xC9,0xCF,0xDB), False),
         ("deploys the fix it can", RED, True), (" — autonomously. Plus an analyst your team can just talk to.",
          RGBColor(0xC9,0xCF,0xDB), False)], y=2.85, w=11)
cy, cw, ch = 4.35, (CW - 0.28)/2, 2.0
card(s, MX, cy, cw, ch, CYAN, "The autonomous detect → remediate loop",
     "Sense → forecast → detect → reason → decide → act → confirm. Runs on its own, every cycle, and takes real action on the store node.", num="Surface 01")
card(s, MX+cw+0.28, cy, cw, ch, VIOLET, "A conversational data analyst",
     "Ask anything in plain English; it plans, queries your warehouse, interprets the result, draws the chart, and remembers the thread.", num="Surface 02")

# 5 · AGENTIC LOOP (HERO)
s = slide()
eyebrow(s, "Agentic AI · Multi-step autonomous action")
title(s, "From ingest to remediation — no human in the middle.", 32)
steps = [
    ("01 · SENSE", "Poll the warehouse", "The scheduler autonomously re-scans ingested POS/ERP/weather every cycle — no human trigger.", CYAN, False),
    ("02 · FORECAST", "Predict the band", "An ML time-series model sets the expected demand range (p05–p95) per store × item × channel.", CYAN, False),
    ("03 · DETECT", "Reality vs forecast", "Actuals scored vs the band → outages, drops, fraud, delivery-partner blackouts flagged.", CYAN, False),
    ("04 · REASON", "Validate & explain", "A multi-judge council kills false alarms; driver attribution names why. The LLM narrates.", CYAN, False),
    ("05 · DECIDE", "Plan the response", "Classify severity; choose the remediation — fixable remotely, and which action?", CYAN, False),
    ("06 · ACT", "Deploy the fix", "Tool use → real action. Dispatch to the store-node agent: restart POS, fail over network, block IP.", RED, True),
    ("07 · CONFIRM", "Close the loop", "The agent executes on the node and reports back; a failed fix self-corrects or escalates.", CYAN, False),
    ("08 · LEARN ↺", "Feed it forward", "Ops confirms actionable/dismissed → tunes thresholds and feeds the next cycle.", CYAN, False),
]
gx, gy = MX, 2.75
cw = (CW - 3*0.2)/4
ch = 1.72
for i, (num, head, body, acc, act) in enumerate(steps):
    x = gx + (i % 4)*(cw+0.2)
    y = gy + (i//4)*(ch+0.2)
    rect(s, x, y, cw, ch, PANEL, RED if act else LINE, line_w=1.5 if act else 0.75)
    box(s, x+0.2, y+0.16, cw-0.4, 0.28, num, 10.5, acc, bold=True, font=MONO, spacing=1.0)
    box(s, x+0.2, y+0.5, cw-0.4, 0.35, head, 14.5, INK, bold=True, line=1.0)
    box(s, x+0.2, y+0.9, cw-0.4, ch-1.0, body, 10.5, MUTED, line=1.28)
box(s, MX, 6.55, 11.5, 0.5, [("↺ CONTINUOUS & AUTONOMOUS   ", CYAN, True),
    ("Planning · reasoning · tool use · action that changes the world", INK, True),
    (" — not a single prompt-response wrapper.", MUTED, False)], 13, MUTED, font=MONO, spacing=0.4, line=1.2)

# 6 · AGENTIC SURFACE 2
s = slide()
eyebrow(s, "Agentic AI · Tool use & self-correction")
title(s, "An analyst that reasons, queries, and self-heals.", 34)
lead(s, "Every question runs a multi-step agent loop over your live warehouse — powered by OpenAI, "
        "guardrailed to read-only SQL over allow-listed tables.", y=2.9, w=11)
cy, cw, ch = 3.95, (CW-0.28)/2, 2.05
card(s, MX, cy, cw, ch, VIOLET, "The loop, per turn",
     "Plan → write SQL (tool) → run on the warehouse → observe rows → self-correct on a DB error → interpret in prose → visualize (bar / line / pie / heatmap). Remembers context across follow-ups.")
card(s, MX+cw+0.28, cy, cw, ch, CYAN, "Live example",
     "“How does weather affect sales?” → joins weather × sales → “26.5% drop on rainy days.”  “So should we staff up?” → “Yes — rain shifts demand to delivery; here's the plan.”")
box(s, MX, 6.25, 11.4, 0.6, "When a query errors, the agent feeds the database error back to itself and rewrites it — "
    "a genuine reason → act → observe cycle, not one prompt.", 14.5, MUTED, line=1.4)

# 7 · TECHNICAL EXECUTION
s = slide()
eyebrow(s, "Technical execution")
title(s, "Built, deployed, and running — on our own hardware.", 32)
flow = [("Ingest", "POS · ERP · weather", False), ("Medallion warehouse", "bronze → silver → gold", True),
        ("ML time-series model", "p05 / p50 / p95", True), ("Anomaly · council · drivers", "detect · explain", True),
        ("Store-node agent", "remediation", True), ("Serve", "dashboard · AI · BI", False)]
fx, fy = MX, 2.85
fw = (CW - 5*0.14)/6
for i, (t, sub, hl) in enumerate(flow):
    x = fx + i*(fw+0.14)
    rect(s, x, fy, fw, 1.0, PANEL, CYAN if hl else LINE, line_w=1.25 if hl else 0.75)
    box(s, x+0.14, fy+0.16, fw-0.28, 0.5, t, 12, INK, bold=True, line=1.05)
    box(s, x+0.14, fy+0.62, fw-0.28, 0.3, sub, 9.5, MUTED, line=1.1)
cy, cw, ch = 4.35, (CW-2*0.28)/3, 1.95
card(s, MX,             cy, cw, ch, CYAN, "Provider-agnostic LLM", "OpenAI today; any OpenAI-compatible or Anthropic endpoint by config. Optional — deterministic fallbacks everywhere.")
card(s, MX+cw+0.28,     cy, cw, ch, CYAN, "Guardrailed & safe", "LLM-written SQL is SELECT-only, single-statement, LIMIT-bound, and can never touch credential tables.")
card(s, MX+2*(cw+0.28), cy, cw, ch, CYAN, "Live on-prem", "No cloud bill — runs on hardware you own, behind TLS, alongside existing systems. Bilingual EN/VI, VND-native.")

# 8 · IMPACT
s = slide()
eyebrow(s, "Impact / usefulness")
title(s, "A same-shift response — at chain scale, no new headcount.", 30)
metrics = [("24–48h → <2h", "Time-to-detect — a >90% cut vs. manual review"),
           ("9.1%", "Store-daily forecast error (MAPE) — beats the ≤10% target"),
           ("100%", "Of anomalies auto-explained with ≥3 root-cause drivers"),
           ("1-click", "Remediation — restart, fail over, block IP, rotate creds"),
           ("250+", "Stores watched continuously — scales without adding analysts"),
           ("₫0", "New cloud spend — runs on infrastructure you already own")]
mx0, my0 = MX, 3.15
mw = (CW - 2*0.4)/3
for i, (nnum, lbl) in enumerate(metrics):
    x = mx0 + (i % 3)*(mw+0.4)
    y = my0 + (i//3)*1.85
    col = RED if i == 0 else (INK if i in (1,) else INK)
    box(s, x, y, mw, 0.8, nnum, 40 if len(nnum) < 8 else 30, RED if i in (0,) else INK, bold=True, spacing=-0.6)
    box(s, x, y+0.85, mw-0.2, 0.9, lbl, 12.5, MUTED, line=1.32)

# 9 · DIFFERENTIATION
s = slide()
eyebrow(s, "What makes it different")
title(s, "Most tools stop at the alert. We deploy the fix.", 32)
cy, cw, ch = 2.9, (CW-0.28)/2, 1.85
card(s, MX,            cy,        cw, ch, RED,    "Security DNA in a restaurant", "The same detect-and-respond discipline that protects enterprises now protects a shift — down to blocking an intruding IP and rotating credentials on a store node.")
card(s, MX+cw+0.28,    cy,        cw, ch, CYAN,   "The loop actually closes", "Detection is table stakes. We dispatch the remediation to the node and confirm it worked — an agent that changes the world, not one that just narrates it.")
card(s, MX,            cy+ch+0.25, cw, ch, VIOLET, "An analyst that draws & self-heals", "Conversational AI that writes its own SQL, renders its own charts and heatmaps, and rewrites a failing query on the fly.")
card(s, MX+cw+0.28,    cy+ch+0.25, cw, ch, AMBER,  "Built for the market, from inside it", "Vietnamese holidays and Tết, delivery partners (Grab / ShopeeFood / Baemin), VND, EN/VI — modeled, not bolted on.")

# 10 · TEAM
s = slide()
eyebrow(s, "The team")
title(s, "We didn't read about this in a case study. We lived it.", 30)
cy, cw, ch = 3.1, (CW-2*0.28)/3, 2.3
card(s, MX,             cy, cw, ch, RED, "Jason Koo", "Co-founder. Former CIO, Vinh Hoan Corp (public, VN) and Director of Engineering, Binary Defense — enterprise systems + threat detection.")
card(s, MX+cw+0.28,     cy, cw, ch, RED, "Kevin Yee", "Co-founder. Former Vietnam Country Manager, Apple; CGO / CMO at CoderSchool & POPS Worldwide — market, growth & P&L.")
card(s, MX+2*(cw+0.28), cy, cw, ch, RED, "Christian Decker", "Co-founder. Engineering & business analysis — turns operational pain into software that ships.")
box(s, MX, 5.75, 11.5, 0.5, "ACROSS OUR CAREERS:  STARBUCKS · APPLE · GE · GEICO  —  LEADERSHIP TO HANDS-ON ENGINEERING",
    12.5, MUTED, bold=True, font=MONO, spacing=1.4)

# 11 · CLOSE
s = slide()
eyebrow(s, "The bottom line")
box(s, MX, 2.2, 11.5, 0.8, "From “what happened yesterday?”", 30, MUTED, bold=True, spacing=-0.4)
box(s, MX, 3.0, 11.5, 0.8, "to “here's what's happening now —", 40, CYAN, bold=True, spacing=-0.5)
box(s, MX, 3.85, 11.5, 0.8, "and here's the fix.”", 40, RED, bold=True, spacing=-0.5)
x = pill(s, MX, 5.35, "Live now · forecaster.secureinsights.ai", GREEN, RGBColor(0x1e,0x4d,0x3a), dot=True)
pill(s, x, 5.35, "Secure Insights AI")
box(s, MX, 6.0, 11, 0.6, "Give us your stores — we'll turn a 24-hour blind spot into a same-shift response.",
    15, MUTED, line=1.4)

# footers on content slides
for sl in list(prs.slides)[1:]:
    footer(sl)

OUT = os.path.join(os.path.dirname(__file__), "..", "secure-insights-pitch.pptx")
prs.save(OUT)
print("saved", os.path.abspath(OUT), "·", len(prs.slides), "slides")
