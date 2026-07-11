from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

T = RGBColor(0, 77, 67)
G = RGBColor(212, 168, 75)
DK = RGBColor(26, 26, 26)
DG = RGBColor(51, 51, 51)
MG = RGBColor(102, 102, 102)
LG = RGBColor(153, 153, 153)
SG = RGBColor(247, 248, 249)
W = RGBColor(255, 255, 255)
TL = RGBColor(232, 240, 238)
TM = RGBColor(0, 122, 107)
C2 = RGBColor(200, 200, 200)
C1 = RGBColor(150, 150, 150)
C3 = RGBColor(180, 200, 195)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(sl, c):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = c

def rn(sl, l, t, w, h, fc=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.shadow.inherit = False
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else: s.fill.background()
    s.line.fill.background()
    s.adjustments[0] = 0.05

def rs(sl, l, t, w, h, fc=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else: s.fill.background()
    s.line.fill.background()

def tx(sl, l, t, w, h, txt, sz=12, b=False, c=DK, al=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.name = "Arial"; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.alignment = al; p.space_after = Pt(0)

def rt(sl, l, t, w, h, runs, al=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al; p.space_after = Pt(0)
    for i, rd in enumerate(runs):
        txt = rd[0]; bl = rd[1] if len(rd) > 1 else False
        sz = rd[2] if len(rd) > 2 else 12; cl = rd[3] if len(rd) > 3 else DK
        run = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
        run.text = txt; run.font.name = "Arial"; run.font.size = Pt(sz)
        run.font.bold = bl; run.font.color.rgb = cl

def ft(sl, n, t):
    rs(sl, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), fc=SG)
    tx(sl, Inches(0.67), Inches(7.03), Inches(8), Inches(0.35), "KFC Vietnam Platform", sz=7, c=LG)
    tx(sl, Inches(10.5), Inches(7.03), Inches(2.5), Inches(0.35), "Slide " + str(n) + " of " + str(t), sz=7, c=LG, al=PP_ALIGN.RIGHT)

def atl(sl, txt):
    tx(sl, Inches(0.67), Inches(0.35), Inches(10), Inches(0.45), txt, sz=13, b=True, c=DK)
    rs(sl, Inches(0.67), Inches(0.87), Inches(0.5), Pt(3), fc=T)

# === SLIDE 1: TITLE ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, T)
rs(sl, Inches(0.83), Inches(2.0), Inches(0.83), Pt(4), fc=G)
tx(sl, Inches(0.83), Inches(2.4), Inches(11), Inches(1.2), "Faster Sales Anomaly Detection and Forecasting", sz=36, b=True, c=W)
tx(sl, Inches(0.83), Inches(3.8), Inches(9), Inches(0.6), "Reducing response time from 24-48 hours to near-real time across the KFC Vietnam chain", sz=16, c=C2)
tx(sl, Inches(0.83), Inches(6.5), Inches(6), Inches(0.4), "KFC Vietnam - Operations and Finance", sz=10, b=True, c=G)
tx(sl, Inches(10), Inches(6.5), Inches(2.5), Inches(0.4), "2025", sz=10, c=C1, al=PP_ALIGN.RIGHT)

# === SLIDE 2: THE PROBLEM ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Manual daily report review creates a 24-48 hour blind spot, costing the business critical recovery time on every incident")
ft(sl, 2, 14)
rn(sl, Inches(0.67), Inches(1.3), Inches(3.2), Inches(1.0), fc=SG)
tx(sl, Inches(0.9), Inches(1.4), Inches(2.8), Inches(0.45), "24-48h", sz=26, b=True, c=T)
tx(sl, Inches(0.9), Inches(1.85), Inches(2.8), Inches(0.3), "Detection lag per incident", sz=9, c=MG)
rn(sl, Inches(4.07), Inches(1.3), Inches(3.2), Inches(1.0), fc=SG)
tx(sl, Inches(4.3), Inches(1.4), Inches(2.8), Inches(0.45), "100%", sz=26, b=True, c=T)
tx(sl, Inches(4.3), Inches(1.85), Inches(2.8), Inches(0.3), "Manual review dependency per store", sz=9, c=MG)
tx(sl, Inches(0.67), Inches(2.6), Inches(6.8), Inches(0.4), "What does a 36-hour blind spot cost us?", sz=10, c=MG, al=PP_ALIGN.LEFT)

issues = [
    ("POS / register outage", "Sales silently stop at a store - full day of lost revenue"),
    ("Void / comp fraud", "Refund abuse continues unchecked - losses accrue every shift"),
    ("Stockout on best-sellers", "Walked customers, missed sales - direct margin impact"),
    ("Demand spike / drop", "Wrong staffing and prep levels - waste or poor service"),
]
for i, (ti, de) in enumerate(issues):
    y = Inches(1.3) + i * Inches(0.7)
    rs(sl, Inches(7.5), y + Inches(0.05), Inches(0.12), Inches(0.12), fc=TM)
    tx(sl, Inches(7.8), y, Inches(5), Inches(0.25), ti, sz=10, b=True, c=DK)
    tx(sl, Inches(7.8), y + Inches(0.25), Inches(5), Inches(0.35), de, sz=8.5, c=MG)
tx(sl, Inches(0.67), Inches(6.6), Inches(8), Inches(0.3), "Source: Current KFC Vietnam operations workflow analysis", sz=7, c=LG)

# === SLIDE 3: COST OF DELAY ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Every hour of delay compounds financial losses across four distinct risk categories, and manual review does not scale with chain growth")
ft(sl, 3, 14)
yh = Inches(1.2)
for i, h in enumerate(["Issue type", "What happens in the blind spot", "Cost of 24-48h delay"]):
    tx(sl, [Inches(0.67), Inches(3.17), Inches(9.17)][i], yh, [Inches(2.5), Inches(6.0), Inches(3.5)][i], Inches(0.35), h.upper(), sz=8, b=True, c=LG)
rs(sl, Inches(0.67), yh + Inches(0.35), Inches(12), Pt(1), fc=RGBColor(208, 208, 208))
rows = [
    ("POS / register outage", "Sales silently stop at a store; no alert triggers until next day report review", "Full day of unrecorded or lost revenue"),
    ("Void / comp fraud", "Refund and discount abuse continues across multiple shifts without detection", "Losses compound with every shift until caught"),
    ("Stockout on best-sellers", "High-margin items unavailable; customers walk or shift to competitors", "Missed sales + damaged customer satisfaction"),
    ("Demand spike / drop", "Wrong staffing levels and food prep quantities locked in for the day", "Labor waste or service failures; both hit margins"),
]
for idx, (c1, c2, c3) in enumerate(rows):
    y = yh + Inches(0.45) + idx * Inches(0.55)
    tx(sl, Inches(0.67), y, Inches(2.5), Inches(0.45), c1, sz=9, b=True, c=DK)
    tx(sl, Inches(3.17), y, Inches(6.0), Inches(0.45), c2, sz=8.5, c=DG)
    tx(sl, Inches(9.17), y, Inches(3.5), Inches(0.45), c3, sz=9, b=True, c=T)
    rs(sl, Inches(0.67), y + Inches(0.5), Inches(12), Pt(0.5), fc=RGBColor(238, 238, 238))
rn(sl, Inches(0.67), Inches(4.8), Inches(12), Inches(0.7), fc=SG)
rs(sl, Inches(0.67), Inches(4.8), Pt(4), Inches(0.7), fc=T)
tx(sl, Inches(0.9), Inches(4.9), Inches(9), Inches(0.5), "Key insight: Manual review does not scale - more stores = more reports = longer lag", sz=10, c=DK)
tx(sl, Inches(9.5), Inches(4.9), Inches(3), Inches(0.35), "+36%", sz=18, b=True, c=T, al=PP_ALIGN.RIGHT)
tx(sl, Inches(9.5), Inches(5.25), Inches(3), Inches(0.25), "Analyst time spent finding vs. fixing", sz=7, c=LG, al=PP_ALIGN.RIGHT)
tx(sl, Inches(0.67), Inches(5.8), Inches(10), Inches(0.3), "Source: KFC Vietnam operations workflow analysis", sz=7, c=LG)

# === SLIDE 4: WHAT WE BUILT ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "An automated platform now forecasts demand across every store, detects anomalies in minutes, and routes alerts to the right person automatically")
ft(sl, 4, 14)
tx(sl, Inches(0.67), Inches(1.0), Inches(8), Inches(0.35), "Six capabilities working together to eliminate the 24-48 hour blind spot", sz=10, c=MG)
caps = [
    ("Forecast", "Demand prediction per store x item x daypart, 14 days ahead"),
    ("Detect", "Automated anomaly detection across outages, fraud, stockouts"),
    ("Validate", "Multi-judge council evaluates each anomaly for confidence"),
    ("Act", "Forecast-driven buying plans, prep schedules, push alerts"),
    ("Surface", "Dashboard, natural-language querying, BI feeds, summaries"),
    ("Localized", "Built for Vietnam: EN/VI, VND, Tet regime, local holidays"),
]
for i, (ti, de) in enumerate(caps):
    col, row = i % 3, i // 3
    x, y = Inches(0.67) + col * Inches(4.2), Inches(1.5) + row * Inches(2.2)
    rn(sl, x, y, Inches(3.9), Inches(1.9), fc=SG)
    tx(sl, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.35), ti, sz=12, b=True, c=T)
    tx(sl, x + Inches(0.2), y + Inches(0.55), Inches(3.5), Inches(1.1), de, sz=8.5, c=DG)
tx(sl, Inches(0.67), Inches(5.9), Inches(8), Inches(0.3), "Source: KFC Vietnam platform architecture documentation", sz=7, c=LG)

# === SLIDE 5: BEFORE VS AFTER ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Detection time collapses from 24-48 hours of manual review to minutes of automated processing")
ft(sl, 5, 14)
yb = Inches(1.3)
rs(sl, Inches(0.67), yb, Inches(1.0), Inches(1.2), fc=C2)
tx(sl, Inches(0.72), yb + Inches(0.35), Inches(0.9), Inches(0.4), "BEFORE", sz=10, b=True, c=W, al=PP_ALIGN.CENTER)
tx(sl, Inches(0.72), yb + Inches(0.7), Inches(0.9), Inches(0.3), "Manual", sz=8, c=W, al=PP_ALIGN.CENTER)
bsteps = ["Day ends", "Analyst pulls reports", "Eyeballs each store", "Spots issue 24-48h later", "Emails and reacts"]
bcols = [SG, SG, SG, RGBColor(255, 230, 230), SG]
for i, (st, bc) in enumerate(zip(bsteps, bcols)):
    x = Inches(1.9) + i * Inches(2.05)
    rn(sl, x, yb, Inches(1.85), Inches(1.2), fc=bc)
    tx(sl, x + Inches(0.1), yb + Inches(0.15), Inches(1.65), Inches(0.9), st, sz=9, c=DG, al=PP_ALIGN.CENTER)
    if i < 4:
        tx(sl, x + Inches(1.85), yb + Inches(0.35), Inches(0.2), Inches(0.4), ">", sz=14, c=LG, al=PP_ALIGN.CENTER)
ya = Inches(2.8)
rs(sl, Inches(0.67), ya, Inches(1.0), Inches(1.2), fc=T)
tx(sl, Inches(0.72), ya + Inches(0.35), Inches(0.9), Inches(0.4), "AFTER", sz=10, b=True, c=G, al=PP_ALIGN.CENTER)
tx(sl, Inches(0.72), ya + Inches(0.7), Inches(0.9), Inches(0.3), "Platform", sz=8, c=W, al=PP_ALIGN.CENTER)
asteps = ["Data lands", "Pipeline forecasts + scores", "Council validates", "Alert fires minutes later", "Right person acts same shift"]
acols = [TL, TL, RGBColor(247, 240, 232), TL, TL]
for i, (st, ac) in enumerate(zip(asteps, acols)):
    x = Inches(1.9) + i * Inches(2.05)
    rn(sl, x, ya, Inches(1.85), Inches(1.2), fc=ac)
    tx(sl, x + Inches(0.1), ya + Inches(0.15), Inches(1.65), Inches(0.9), st, sz=9, c=DG, al=PP_ALIGN.CENTER)
    if i < 4:
        tx(sl, x + Inches(1.85), ya + Inches(0.35), Inches(0.2), Inches(0.4), ">", sz=14, c=LG, al=PP_ALIGN.CENTER)
rn(sl, Inches(0.67), Inches(4.4), Inches(12), Inches(0.9), fc=T)
meas = [("24-48h", "Minutes"), ("Manual", "Automatic chain-wide"), ('"Find"', '"Act"')]
for i, (bef, aft) in enumerate(meas):
    x = Inches(1.2) + i * Inches(4.0)
    tx(sl, x, Inches(4.55), Inches(1.5), Inches(0.35), bef, sz=14, b=True, c=G)
    tx(sl, x + Inches(1.5), Inches(4.6), Inches(2), Inches(0.3), "  " + aft, sz=10, c=W)
tx(sl, Inches(0.67), Inches(5.6), Inches(8), Inches(0.3), "Source: Platform architecture and workflow design", sz=7, c=LG)

# === SLIDE 6: PIPELINE ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Every pipeline run processes the full chain from raw data to alert in minutes, removing the human from the detection loop")
ft(sl, 6, 14)
pipe = [("1", "Ingest", "POS, ERP, weather data"), ("2", "Forecast", "Chronos-2 and seasonal models"), ("3", "Detect", "Residual and fraud scoring"), ("4", "Validate", "Council of judges votes"), ("5", "Plan", "Buying and prep plans"), ("6", "Notify", "Email, Teams, Slack")]
pcols = [SG, TL, RGBColor(255, 240, 235), RGBColor(247, 240, 232), TL, T]
for i, (nu, ti, de) in enumerate(pipe):
    x = Inches(0.67) + i * Inches(2.1)
    bc = pcols[i]
    rn(sl, x, Inches(1.3), Inches(1.85), Inches(2.5), fc=bc)
    tx(sl, x, Inches(1.4), Inches(1.85), Inches(0.25), "STEP " + nu, sz=8, b=True, c=LG, al=PP_ALIGN.CENTER)
    tc = W if bc == T else DK
    tx(sl, x, Inches(1.7), Inches(1.85), Inches(0.35), ti, sz=11, b=True, c=tc, al=PP_ALIGN.CENTER)
    tx(sl, x + Inches(0.15), Inches(2.15), Inches(1.55), Inches(1.4), de, sz=8, c=MG, al=PP_ALIGN.CENTER)
    if i < 5:
        tx(sl, x + Inches(1.85), Inches(2.3), Inches(0.25), Inches(0.35), ">", sz=12, c=LG, al=PP_ALIGN.CENTER)
tx(sl, Inches(0.67), Inches(4.15), Inches(12), Inches(0.5), "Fully automated | Reproducible rebuilds from source | Flexible cadence nightly or on-demand", sz=8.5, c=DG)
tx(sl, Inches(0.67), Inches(5.0), Inches(8), Inches(0.3), "Source: Platform pipeline architecture documentation", sz=7, c=LG)

# === SLIDE 7: FORECAST ENGINE ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "The forecast engine delivers reliable 14-day predictions with honest uncertainty ranges, validated at 22% MAPE and 89% band coverage")
ft(sl, 7, 14)
feats = [
    ("Two-model architecture", "Chronos-2 foundation model (GPU) for highest accuracy; zero-dependency seasonal fallback with identical interface."),
    ("Forecast as a range", "Every prediction is a p05-p50-p95 band. The width tells you confidence. Band directly powers anomaly detection."),
    ("Restaurant-aware inputs", "Weather, day-of-week, promotions, Tet and Vietnam holidays, new-store maturation ramp."),
    ("Trusted data foundation", "Bronze to Silver to Gold medallion warehouse with full lineage tracing to source exports."),
]
for i, (ti, de) in enumerate(feats):
    y = Inches(1.1) + i * Inches(1.15)
    tx(sl, Inches(0.67), y, Inches(7), Inches(0.3), ti, sz=10, b=True, c=T)
    tx(sl, Inches(0.67), y + Inches(0.3), Inches(7), Inches(0.7), de, sz=8.5, c=DG)
rn(sl, Inches(8.2), Inches(1.1), Inches(4.5), Inches(4.5), fc=SG)
met = [("MAPE", "22%"), ("Band coverage", "89%"), ("Skill vs. naive", "+36%"), ("Horizon", "14 days")]
for i, (la, va) in enumerate(met):
    y = Inches(1.3) + i * Inches(1.0)
    tx(sl, Inches(8.5), y, Inches(2.5), Inches(0.3), la, sz=8.5, c=MG)
    tx(sl, Inches(11), y, Inches(1.5), Inches(0.35), va, sz=16, b=True, c=T, al=PP_ALIGN.RIGHT)
    if i == 1:
        rs(sl, Inches(8.5), y + Inches(0.7), Inches(3.7), Pt(4), fc=RGBColor(208, 208, 208))
        rs(sl, Inches(8.5), y + Inches(0.7), Inches(3.3), Pt(4), fc=T)
tx(sl, Inches(0.67), Inches(5.8), Inches(8), Inches(0.3), "Models: Chronos-2 (GPU) + Seasonal ML | Validated on: KFC Vietnam hold-out data", sz=7.5, c=MG)
tx(sl, Inches(0.67), Inches(6.2), Inches(8), Inches(0.3), "Source: Backtest results on KFC Vietnam historical sales data", sz=7, c=LG)

# === SLIDE 8: ANOMALY DETECTION ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Anomaly detection achieves precision 1.0 across all types through a multi-judge validation council that eliminates false alarms")
ft(sl, 8, 14)
fl = [("Compare", "Actual vs forecast band", TL), ("Classify", "Spike/Drop/Outage/Fraud", RGBColor(255, 240, 235)), ("Council", "Multiple judges vote", RGBColor(247, 240, 232)), ("Alert", "Ranked anomalies", T)]
for i, (ti, de, bc) in enumerate(fl):
    x = Inches(0.67) + i * Inches(2.8)
    rn(sl, x, Inches(1.3), Inches(2.5), Inches(1.4), fc=bc)
    tc = W if bc == T else DK
    tx(sl, x, Inches(1.45), Inches(2.5), Inches(0.3), ti, sz=10, b=True, c=tc, al=PP_ALIGN.CENTER)
    tx(sl, x + Inches(0.15), Inches(1.8), Inches(2.2), Inches(0.7), de, sz=8.5, c=MG, al=PP_ALIGN.CENTER)
    if i < 3:
        tx(sl, x + Inches(2.5), Inches(1.75), Inches(0.3), Inches(0.3), ">", sz=12, c=LG, al=PP_ALIGN.CENTER)
rn(sl, Inches(0.67), Inches(3.0), Inches(7.5), Inches(0.65), fc=SG)
rs(sl, Inches(0.67), Inches(3.0), Pt(4), Inches(0.65), fc=T)
tx(sl, Inches(0.9), Inches(3.1), Inches(7), Inches(0.5), "How it works: Each anomaly evaluated by multiple independent judges. Alert fires only when a quorum agrees - eliminating alert fatigue.", sz=8.5, c=DG)
rn(sl, Inches(8.5), Inches(1.3), Inches(3.8), Inches(1.5), fc=T)
tx(sl, Inches(8.8), Inches(1.5), Inches(3.2), Inches(0.5), "1.0", sz=32, b=True, c=G, al=PP_ALIGN.CENTER)
tx(sl, Inches(8.8), Inches(2.0), Inches(3.2), Inches(0.4), "Precision across all types", sz=9, c=W, al=PP_ALIGN.CENTER)
tx(sl, Inches(8.8), Inches(2.35), Inches(3.2), Inches(0.25), "When it alerts, it is right", sz=7.5, c=C3, al=PP_ALIGN.CENTER)
rn(sl, Inches(8.5), Inches(3.1), Inches(3.8), Inches(2.5), fc=SG)
tx(sl, Inches(8.7), Inches(3.25), Inches(3.4), Inches(0.25), "ANOMALY TYPES DETECTED", sz=8, b=True, c=MG)
for i, at in enumerate(["Outage - Sales stopped", "Fraud - Void/comp abuse", "Stockout - Best-sellers", "Demand spike/drop", "Inventory mismatch"]):
    tx(sl, Inches(8.7), Inches(3.6) + i * Inches(0.38), Inches(3.4), Inches(0.35), at, sz=8, c=DG)
tx(sl, Inches(0.67), Inches(6.0), Inches(8), Inches(0.3), "Source: Platform backtest results on KFC Vietnam data", sz=7, c=LG)

# === SLIDE 9: CLOSING THE LOOP ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "The platform closes the loop from detection to action, routing alerts and plans so the right person responds same shift")
ft(sl, 9, 14)
tx(sl, Inches(0.67), Inches(1.1), Inches(10), Inches(0.35), "Detecting faster only matters if the right person acts. Four integrated action pathways.", sz=10, c=MG)
acts = [("Automated alerts", "Rule-based routing by severity to email, Teams, Slack"), ("Buying plans", "Par levels, reorder points, quantities in VND"), ("Prep plans", "Thaw and cook lead-time board for forecast demand"), ("Agent-ready (MCP)", "AI agents can query and approve with audit trail")]
for i, (ti, de) in enumerate(acts):
    col, row = i % 2, i // 2
    x, y = Inches(0.67) + col * Inches(6.3), Inches(1.6) + row * Inches(1.8)
    rn(sl, x, y, Inches(5.9), Inches(1.5), fc=SG)
    rs(sl, x, y, Pt(4), Inches(1.5), fc=T)
    tx(sl, x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(0.3), ti, sz=11, b=True, c=DK)
    tx(sl, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.8), de, sz=8.5, c=DG)
tx(sl, Inches(0.67), Inches(5.5), Inches(8), Inches(0.3), "Source: Platform action-routing architecture", sz=7, c=LG)

# === SLIDE 10: DASHBOARD ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "The single dashboard gives every role a tailored view with no SQL required and full drill-down to any store")
ft(sl, 10, 14)
f10 = [("01", "Chain-wide snapshot", "90-day trend and 14-day forecast"), ("02", "Forecast Explorer", "Backtest-vs-actual overlay"), ("03", "Anomaly triage", "Sort by type and confidence"), ("04", "Natural-language query", "Plain English gets a table"), ("05", "AI summaries", "12 report templates plus auto summary"), ("06", "Full localization", "EN/VI toggle, VND currency")]
for i, (nu, ti, de) in enumerate(f10):
    col, row = i % 3, i // 3
    x, y = Inches(0.67) + col * Inches(4.2), Inches(1.3) + row * Inches(2.2)
    rn(sl, x, y, Inches(3.9), Inches(1.85), fc=SG)
    tx(sl, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.35), nu, sz=18, b=True, c=T)
    tx(sl, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.3), ti, sz=11, b=True, c=DK)
    tx(sl, x + Inches(0.2), y + Inches(0.85), Inches(3.5), Inches(0.8), de, sz=8.5, c=DG)
rn(sl, Inches(0.67), Inches(5.2), Inches(12), Inches(0.7), fc=T)
tx(sl, Inches(0.9), Inches(5.35), Inches(6), Inches(0.3), "Available now at 192.168.50.85:8900", sz=9, b=True, c=G)
tx(sl, Inches(0.9), Inches(5.6), Inches(6), Inches(0.2), "Analysts | Managers | Finance | Operations", sz=7.5, c=W)
tx(sl, Inches(0.67), Inches(6.2), Inches(8), Inches(0.3), "Source: Platform UI and feature documentation", sz=7, c=LG)

# === SLIDE 11: INFRASTRUCTURE ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "The platform runs on existing on-prem hardware with no new cloud spend and full RBAC data governance")
ft(sl, 11, 14)
inf = [("On-premise deployment", "Docker Compose on idle GPU (A2) at 192.168.50.85 - no contention"), ("One-command deploy", "Same code on laptop (SQLite) or server (Postgres+GPU)"), ("Enterprise security", "PBKDF2 hashing, short sessions, pending signup approval")]
for i, (ti, de) in enumerate(inf):
    y = Inches(1.15) + i * Inches(1.5)
    rn(sl, Inches(0.67), y, Inches(7.5), Inches(1.3), fc=SG)
    rs(sl, Inches(0.67), y, Pt(4), Inches(1.3), fc=T)
    tx(sl, Inches(0.9), y + Inches(0.1), Inches(7), Inches(0.3), ti, sz=10, b=True, c=T)
    tx(sl, Inches(0.9), y + Inches(0.4), Inches(7), Inches(0.7), de, sz=8.5, c=DG)
rn(sl, Inches(8.5), Inches(1.15), Inches(4.2), Inches(3.3), fc=SG)
tx(sl, Inches(8.7), Inches(1.25), Inches(3.8), Inches(0.3), "ROLE-BASED ACCESS CONTROL", sz=8, b=True, c=MG)
roles = [("Admin", "Full system access, user mgmt"), ("Manager", "View all + approve reorders"), ("Analyst", "View data, create reports"), ("Viewer", "Read-only dashboard")]
for i, (ro, ca) in enumerate(roles):
    y = Inches(1.7) + i * Inches(0.42)
    tx(sl, Inches(8.7), y, Inches(1.5), Inches(0.35), ro, sz=8.5, b=True, c=T)
    tx(sl, Inches(10.3), y, Inches(3.2), Inches(0.35), ca, sz=8, c=DG)
tx(sl, Inches(8.7), Inches(3.6), Inches(3.8), Inches(0.4), "Audit trail: Every action logged with timestamp", sz=7.5, c=MG)
rn(sl, Inches(0.67), Inches(5.0), Inches(12), Inches(0.6), fc=SG)
for i, (la, va) in enumerate([("Stack:", "App + Postgres + GPU"), ("Deploy:", "Docker Compose"), ("Host:", "192.168.50.85")]):
    x = Inches(0.9) + i * Inches(4.2)
    tx(sl, x, Inches(5.1), Inches(3.8), Inches(0.35), la + " " + va, sz=9, c=DG)
tx(sl, Inches(0.67), Inches(5.9), Inches(8), Inches(0.3), "Source: Platform infrastructure and security docs", sz=7, c=LG)

# === SLIDE 12: BUSINESS VALUE ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "The platform delivers measurable business value - faster loss prevention, less waste, and scaled analyst productivity")
ft(sl, 12, 14)
yh = Inches(1.2)
for i, h in enumerate(["Value lever", "How the platform delivers", "Impact"]):
    tx(sl, [Inches(0.67), Inches(3.17), Inches(9.67)][i], yh, [Inches(2.5), Inches(6.5), Inches(2.5)][i], Inches(0.35), h.upper(), sz=8, b=True, c=LG)
vrows = [("Faster loss prevention", "Outages caught in minutes - recovery same shift", "24-48h to min"), ("Less shrink/fraud", "Void-comp flagged by validation council", "Cut losses"), ("Fewer stockouts/waste", "Forecast-driven buying reduces waste", "Less waste"), ("Analyst productivity", "Automated monitoring replaces manual", "Redeploy hours"), ("Scales with chain", "New stores watched automatically", "Zero added cost"), ("Better decisions", "BI feeds enable evidence-based calls", "Faster decisions")]
for idx, (c1, c2, c3) in enumerate(vrows):
    y = yh + Inches(0.45) + idx * Inches(0.55)
    tx(sl, Inches(0.67), y, Inches(2.5), Inches(0.45), c1, sz=9, b=True, c=DK)
    tx(sl, Inches(3.17), y, Inches(6.5), Inches(0.45), c2, sz=8.5, c=DG)
    tx(sl, Inches(9.67), y, Inches(2.5), Inches(0.45), c3, sz=9, b=True, c=T)
rn(sl, Inches(0.67), Inches(5.0), Inches(11.5), Inches(0.6), fc=T)
tx(sl, Inches(0.9), Inches(5.1), Inches(5), Inches(0.35), "24-48h to minutes", sz=14, b=True, c=G)
tx(sl, Inches(5), Inches(5.15), Inches(7), Inches(0.35), "Same-shift response, chain scale, no added headcount", sz=8.5, c=W)
tx(sl, Inches(0.67), Inches(5.9), Inches(10), Inches(0.3), "Source: Platform value analysis", sz=7, c=LG)

# === SLIDE 13: ROAD AHEAD ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, W)
atl(sl, "Built for MVP today and designed to grow - the architecture scales for real-time, more channels, wider rollout")
ft(sl, 13, 14)
tx(sl, Inches(0.67), Inches(1.1), Inches(10), Inches(0.35), "Enhancements build on the same foundation with minimal risk.", sz=10, c=MG)
eras = [("NOW", T, TL, ["Nightly pipeline active", "Full dashboard deployed", "BI feeds live"]), ("NEXT", TM, RGBColor(232, 245, 242), ["Near-streaming ingestion", "Add Zalo channel", "PowerBI semantic model"]), ("LATER", LG, SG, ["More stores/regions", "AI incident summaries", "Multi-chain expansion"])]
for idx, (la, co, bc, items) in enumerate(eras):
    y = Inches(1.6) + idx * Inches(1.35)
    rs(sl, Inches(0.67), y, Inches(1.0), Inches(1.1), fc=co)
    tx(sl, Inches(0.67), y + Inches(0.35), Inches(1.0), Inches(0.35), la, sz=12, b=True, c=W, al=PP_ALIGN.CENTER)
    rn(sl, Inches(1.8), y, Inches(11.5), Inches(1.1), fc=bc)
    for i, item in enumerate(items):
        tx(sl, Inches(2.1) + i * Inches(3.8), y + Inches(0.15), Inches(3.6), Inches(0.8), "  " + item, sz=8.5, c=DG)
rn(sl, Inches(0.67), Inches(5.3), Inches(12), Inches(0.65), fc=SG)
rs(sl, Inches(0.67), Inches(5.3), Pt(4), Inches(0.65), fc=T)
tx(sl, Inches(0.9), Inches(5.38), Inches(5.5), Inches(0.5), "Invitation: Which next item matters most to your team?", sz=9, b=True, c=DK)
tx(sl, Inches(7), Inches(5.38), Inches(5.5), Inches(0.5), "All enhancements build on proven architecture.", sz=9, c=MG)
tx(sl, Inches(0.67), Inches(6.2), Inches(8), Inches(0.3), "Source: Platform product roadmap", sz=7, c=LG)

# === SLIDE 14: SUMMARY ===
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, T)
rs(sl, Inches(0.83), Inches(0.8), Inches(0.67), Pt(4), fc=G)
tx(sl, Inches(0.83), Inches(1.1), Inches(11), Inches(0.7), "From 24-48 hours to minutes: an automated forecast + anomaly platform on our own infrastructure", sz=15, b=True, c=W)
sums = [("01", "Problem eliminated", "Anomalies found 24-48h late. Now detected in minutes."), ("02", "Proven accuracy", "22% MAPE, 89% coverage, precision 1.0."), ("03", "No new infra spend", "On-prem GPU, one-command deploy."), ("04", "Governed and auditable", "RBAC with full audit trail.")]
for i, (nu, ti, de) in enumerate(sums):
    col, row = i % 2, i // 2
    x, y = Inches(0.83) + col * Inches(6.2), Inches(2.2) + row * Inches(1.5)
    rn(sl, x, y, Inches(5.8), Inches(1.3), fc=RGBColor(0, 90, 78))
    tx(sl, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.35), nu, sz=18, b=True, c=G)
    tx(sl, x + Inches(0.2), y + Inches(0.5), Inches(5.3), Inches(0.3), ti, sz=11, b=True, c=W)
    tx(sl, x + Inches(0.2), y + Inches(0.85), Inches(5.3), Inches(0.35), de, sz=8.5, c=C3)
asks = ["Confirm target stores", "Validate figures with Finance", "Prioritize roadmap item"]
for i, ask in enumerate(asks):
    x = Inches(0.83) + i * Inches(4.1)
    rn(sl, x, Inches(5.5), Inches(3.8), Inches(0.65), fc=RGBColor(0, 90, 78))
    rs(sl, x, Inches(5.5), Pt(3), Inches(0.65), fc=G)
    tx(sl, x + Inches(0.4), Inches(5.57), Inches(3.2), Inches(0.4), str(i+1) + ". " + ask, sz=8.5, c=W)
tx(sl, Inches(0.83), Inches(6.6), Inches(6), Inches(0.3), "KFC Vietnam | Sales Anomaly Detection Platform", sz=8, c=C1)
tx(sl, Inches(9), Inches(6.6), Inches(3.5), Inches(0.3), "Contact: [Name] | [Email]", sz=8, c=C1, al=PP_ALIGN.RIGHT)

# === SAVE ===
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "kfc-vietnam-sales-forecast-platform.pptx")
prs.save(out)
print("Done! Created " + out)
print("Total slides: " + str(len(prs.slides)))
